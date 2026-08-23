"""Study tools built on top of the Moodle REST layer.

These are the read-side features that make the bot useful beyond deadlines:

- ``my_progress``     one call, percent complete across every unit
- ``whats_left``      the activities in a unit still not completed
- ``read_material``   downloads the lecturer's own slides or notes and pulls
                      the text out, so the agent explains the real material
                      instead of improvising from memory

``read_material`` is the honest version of "explain this topic to me". The
model is grounded in the file the lecturer actually uploaded, and the reply can
cite the filename. Nothing here writes to Moodle.

Extraction dependencies (python-pptx, python-docx, pypdf) are imported lazily
inside each branch. If one is missing the tool returns a clear message and the
download link, rather than taking the whole service down mid-demo.
"""

from __future__ import annotations

import io
import logging
import re
from typing import Any

import requests

from . import moodle

logger = logging.getLogger(__name__)

# Slide decks in this course reach 8.5 MB. Read generously, but never let one
# file blow the container's memory.
MAX_READ_BYTES = 20 * 1024 * 1024
# Enough text for the model to explain a topic properly, short enough to stay
# cheap on every turn.
MAX_TEXT_CHARS = 8000
MAX_PDF_PAGES = 40

_STOPWORDS = frozenset(
    {
        "the", "a", "an", "and", "or", "of", "in", "on", "for", "to", "is", "are",
        "what", "about", "me", "my", "our", "this", "that", "notes", "note",
        "slides", "slide", "file", "files", "read", "explain", "tell", "unit",
        "topic", "lecture", "send", "give", "please", "it", "i",
    }
)

READABLE_SUFFIXES = (".pptx", ".docx", ".pdf", ".txt", ".md", ".csv")


def _words(text: str) -> set[str]:
    return {
        word
        for word in re.split(r"[^a-z0-9]+", (text or "").lower())
        if len(word) > 2 and word not in _STOPWORDS
    }


def _download(file_url: str, user_key: str) -> bytes:
    """Fetches a Moodle file server-side, with the student's own token."""
    joiner = "&" if "?" in file_url else "?"
    token = moodle._token_for(user_key or None)
    response = requests.get(
        file_url + joiner + "token=" + token,
        timeout=moodle.FILE_TIMEOUT,
        stream=True,
    )
    response.raise_for_status()

    chunks: list[bytes] = []
    total = 0
    for chunk in response.iter_content(65536):
        total += len(chunk)
        if total > MAX_READ_BYTES:
            break
        chunks.append(chunk)
    return b"".join(chunks)


def _extract(data: bytes, filename: str) -> str:
    """Pulls plain text out of a slide deck, document, or PDF.

    Raises:
        ValueError: if the format is not readable, or its library is missing.
    """
    lower = filename.lower()
    buffer = io.BytesIO(data)

    if lower.endswith(".pptx"):
        try:
            from pptx import Presentation
        except ImportError:
            raise ValueError("pptx reading isn't available on this deployment.")
        parts: list[str] = []
        for index, slide in enumerate(Presentation(buffer).slides, start=1):
            lines = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
            ]
            if lines:
                parts.append(f"[Slide {index}] " + " | ".join(lines))
        return "\n".join(parts)

    if lower.endswith(".docx"):
        try:
            import docx
        except ImportError:
            raise ValueError("docx reading isn't available on this deployment.")
        document = docx.Document(buffer)
        parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    if lower.endswith(".pdf"):
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ValueError("pdf reading isn't available on this deployment.")
        reader = PdfReader(buffer)
        parts = []
        for page in reader.pages[:MAX_PDF_PAGES]:
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(text)
        return "\n".join(parts)

    if lower.endswith((".txt", ".md", ".csv")):
        return data.decode("utf-8", errors="replace")

    raise ValueError(
        f"I can't read {filename} as text. I can read pptx, docx, pdf and txt."
    )


# ---------------------------------------------------------------------------
# Tools
# ---------------------------------------------------------------------------


def my_progress(tool_context: Any = None) -> dict:
    """Reports how far this student has got in each of their units.

    Use this for "how am I doing", "which unit am I behind in", or as a quick
    overview before drilling into one unit. This is Moodle's own activity
    completion percentage, not a grade.

    Returns:
        dict: status, plus each unit with its completion percentage.
    """
    user_key = moodle._key(tool_context)
    try:
        courses = moodle._call(
            "core_enrol_get_users_courses",
            {"userid": moodle._user_id(user_key or None)},
            user_key=user_key or None,
        )
    except moodle.MoodleError as error:
        return moodle._fail(error)

    units = []
    for course in courses:
        if course.get("hidden"):
            continue
        progress = course.get("progress")
        units.append(
            {
                "course_id": course.get("id"),
                "code": course.get("shortname"),
                "percent_complete": round(progress, 1)
                if isinstance(progress, (int, float))
                else None,
            }
        )

    # Weakest unit first: that is the one worth talking about.
    units.sort(key=lambda u: (u["percent_complete"] is None, u["percent_complete"] or 0))
    tracked = [u["percent_complete"] for u in units if u["percent_complete"] is not None]
    return {
        "status": "success",
        "units": units,
        "average_percent": round(sum(tracked) / len(tracked), 1) if tracked else None,
        "note": "Completion tracking, not marks. Some units track nothing at all.",
    }


def whats_left(course_id: int, tool_context: Any = None) -> dict:
    """Lists what is still not completed in one unit.

    Splits the result into items the student can tick themselves and items
    Moodle completes for them (by viewing, submitting, or attempting). Useful
    for "what have I not done in mobile programming".

    Args:
        course_id (int): Moodle course id, from list_my_courses.

    Returns:
        dict: status, plus outstanding items grouped by how they complete.
    """
    user_key = moodle._key(tool_context)
    try:
        payload = moodle._call(
            "core_completion_get_activities_completion_status",
            {"courseid": course_id, "userid": moodle._user_id(user_key or None)},
            user_key=user_key or None,
        )
    except moodle.MoodleError as error:
        return moodle._fail(error)

    tickable, automatic, done = [], [], 0
    for item in payload.get("statuses", []):
        if item.get("state"):
            done += 1
            continue
        entry = {"cmid": item.get("cmid"), "type": item.get("modname")}
        if item.get("tracking") == 1:
            tickable.append(entry)
        else:
            automatic.append(entry)

    return {
        "status": "success",
        "course_id": course_id,
        "completed_count": done,
        "you_can_tick": tickable[:25],
        "moodle_completes_these": automatic[:25],
        "note": (
            "Items under moodle_completes_these finish by viewing, submitting "
            "or attempting them. They cannot be ticked by hand."
        ),
    }


def read_material(course_id: int, topic: str, tool_context: Any = None) -> dict:
    """Reads the lecturer's actual slides or notes on a topic, and returns the text.

    Use this whenever a student asks what a topic is about, asks you to explain
    or summarise a week's material, or asks a question that should be answered
    from the unit's own content. Answer from the text this returns, and name the
    file you used. Do not answer from memory when this tool can give you the
    real material.

    Args:
        course_id (int): Moodle course id, from list_my_courses.
        topic (str): What to look for, for example "recyclerview" or
            "testing" or "content providers".

    Returns:
        dict: status, the filename used, its text, and other close matches.
    """
    user_key = moodle._key(tool_context)
    try:
        sections = moodle._call(
            "core_course_get_contents",
            {"courseid": course_id},
            user_key=user_key or None,
        )
    except moodle.MoodleError as error:
        return moodle._fail(error)

    wanted = _words(topic)
    candidates = []
    for section in sections:
        if not section.get("uservisible", True):
            continue
        for module in section.get("modules", []) or []:
            if not module.get("uservisible", True):
                continue
            for item in module.get("contents", []) or []:
                filename = item.get("filename") or ""
                if item.get("type") != "file" or not item.get("fileurl"):
                    continue
                if not filename.lower().endswith(READABLE_SUFFIXES):
                    continue
                haystack = _words(f"{filename} {module.get('name', '')} {section.get('name', '')}")
                candidates.append(
                    {
                        "filename": filename,
                        "fileurl": item["fileurl"],
                        "mimetype": item.get("mimetype", ""),
                        "section": section.get("name"),
                        "score": len(wanted & haystack),
                    }
                )

    if not candidates:
        return {
            "status": "not_found",
            "error_message": "This unit has no readable slides or documents.",
        }

    # Best keyword match, and on a tie prefer the most recent section, which is
    # what a student asking "what are we doing now" almost always means.
    ranked = sorted(
        enumerate(candidates), key=lambda pair: (-pair[1]["score"], -pair[0])
    )
    best = ranked[0][1]

    if best["score"] == 0:
        return {
            "status": "not_found",
            "error_message": (
                f"Nothing in this unit matches '{topic}'. Ask the student to "
                "name the topic differently, or use whats_new_in_unit to list "
                "what is actually there."
            ),
            "available": [c["filename"] for c in candidates[-12:]],
        }

    try:
        data = _download(best["fileurl"], user_key)
        text = _extract(data, best["filename"]).strip()
    except ValueError as error:
        media_id = moodle.register_file(
            best["fileurl"], best["filename"], best["mimetype"], user_key=user_key
        )
        return {
            "status": "unreadable",
            "filename": best["filename"],
            "error_message": str(error),
            "link": moodle.media_path(media_id),
        }
    except Exception:
        logger.exception("failed to read %s", best["filename"])
        return {
            "status": "error",
            "error_message": f"Could not download {best['filename']} from Moodle.",
        }

    if not text:
        media_id = moodle.register_file(
            best["fileurl"], best["filename"], best["mimetype"], user_key=user_key
        )
        return {
            "status": "unreadable",
            "filename": best["filename"],
            "error_message": (
                "That file holds no extractable text - it is probably images or "
                "scans. Send the student the link instead."
            ),
            "link": moodle.media_path(media_id),
        }

    truncated = len(text) > MAX_TEXT_CHARS
    media_id = moodle.register_file(
        best["fileurl"], best["filename"], best["mimetype"], user_key=user_key
    )
    return {
        "status": "success",
        "filename": best["filename"],
        "section": best["section"],
        "truncated": truncated,
        "text": text[:MAX_TEXT_CHARS],
        "link": moodle.media_path(media_id),
        "other_matches": [c["filename"] for _, c in ranked[1:4] if c["score"] > 0],
    }


STUDY_TOOLS = [my_progress, whats_left, read_material]
